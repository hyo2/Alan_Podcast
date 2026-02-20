"""
Metadata Generator Node (V3 - ONNX Recognition + Gemini Fallback)
=================================================================

변경사항:
- EasyOCR/PaddleOCR 완전 제거 (메모리 OOM 해결)
- ONNX Runtime 기반의 초경량 Recognition 아키텍처 도입
- 텍스트 라인 검출(Heuristic Crop) -> ONNX 추론
- 모델 부재 시 Gemini Vision으로 자동 Fallback

"""

import os
import json
import tempfile
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple
from datetime import datetime
import traceback
import io
import logging
import sys
import math

# ✅ 경량화된 라이브러리 임포트
import cv2
import numpy as np
from PIL import Image

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

# 기존 모듈 임포트
from .document_converter_node import DocumentConverterNode
from .improved_hybrid_filter import (
    ImprovedHybridFilterPipeline,
    UniversalImageExtractor,
    ImageMetadata,
    get_global_model,
    gemini_ocr_image_bytes
)
from vertexai.generative_models import Part
from pypdfium2 import PdfDocument

logger = logging.getLogger(__name__)

def _log(*args, level: str | None = None, exc_info: bool = False, end: str = '\n', flush: bool = False) -> None:
    msg = " ".join(str(a) for a in args).rstrip() if args else ""
    if end != "\n" or flush:
        print(msg, end=end, flush=flush)
        return
    lvl = (level or "DEBUG").upper()
    if lvl == "INFO": logger.info(msg, exc_info=exc_info)
    elif lvl in ("WARN", "WARNING"): logger.warning(msg, exc_info=exc_info)
    elif lvl == "ERROR": logger.error(msg, exc_info=exc_info)
    else: logger.debug(msg, exc_info=exc_info)

# ==========================================
# 🔧 RapidOCR Wrapper
# ==========================================
_rapid_ocr_engine = None

def get_rapid_ocr():
    global _rapid_ocr_engine
    if _rapid_ocr_engine is not None:
        return _rapid_ocr_engine
    try:
        from rapidocr_onnxruntime import RapidOCR
        base_dir = Path(__file__).parent.parent.parent / "ocr_model"
        det_path = base_dir / "det.onnx"
        rec_path = base_dir / "rec.onnx"
        dict_path = base_dir / "dict.txt"

        if not det_path.exists() or not rec_path.exists():
            _log(f"⚠️ OCR 모델 파일 없음 ({base_dir}) -> Gemini Fallback", level="WARNING")
            return None

        _rapid_ocr_engine = RapidOCR(
            det_model_path=str(det_path),
            rec_model_path=str(rec_path),
            rec_keys_path=str(dict_path),
        )
        _log("✅ RapidOCR 초기화 완료", level="INFO")
        return _rapid_ocr_engine
    except Exception as e:
        _log(f"⚠️ RapidOCR 초기화 실패: {e}", level="WARNING")
        return None

# ==========================================
# 🔧 Main Class
# ==========================================
class TextExtractor:
    """
    PDF에서 페이지별 텍스트 추출 + 마커 삽입
    V3: pdfplumber + ONNX(Recognition) + Gemini Fallback
    """

    def __init__(self):
        if not PDFPLUMBER_AVAILABLE:
            raise ImportError("pdfplumber가 필요합니다")

        self.ocr_enabled = True
        self.min_text_length = 100
        self.gemini_ocr_fallback = os.getenv('GEMINI_OCR_FALLBACK', 'true').lower() in ('1','true','yes','y')
        self.gemini_ocr_max_sample_pages = int(os.getenv('GEMINI_OCR_MAX_SAMPLE_PAGES', '10'))
        self._gemini_ocr_used_pages = 0
        self._gemini_ocr_skipped_pages = 0

        # RapidOCR 초기화 시도
        self._ocr = get_rapid_ocr()

    def _perform_ocr_on_page(self, pdf_path: str, page_number: int) -> Tuple[str, Optional[Image.Image]]:
        """
        페이지에 OCR 수행
        전략: ONNX (1순위) -> 실패/결과부족 -> Gemini (2순위)
        """
        pil_img = None
        try:
            pdf = PdfDocument(pdf_path)
            page = pdf[page_number - 1]
            bitmap = page.render(scale=2.0)
            pil_img = bitmap.to_pil()

            max_dim = 1024
            if max(pil_img.size) > max_dim:
                pil_img.thumbnail((max_dim, max_dim), Image.LANCZOS)

            if self._ocr is None:
                return "", pil_img

            import numpy as np
            img_np = np.array(pil_img)
            result, elapsed = self._ocr(img_np)

            if not result:
                _log(f"⚠️ RapidOCR 결과 없음 (page {page_number})", level="WARNING")
                return "", pil_img

            texts = [line[1] for line in result if line[1] and line[1].strip()]
            extracted_text = "\n".join(texts)
            _log(f"🧩 RapidOCR 결과: {len(extracted_text)}자 (page {page_number})", level="DEBUG")
            return extracted_text, pil_img

        except Exception as e:
            _log(f"❌ OCR 처리 중 오류 (page {page_number}): {e}", level="ERROR")
            return "", pil_img

    def _calculate_sample_pages(self, total_pages: int, max_samples: int) -> List[int]:
        # (기존 코드 유지)
        if total_pages <= max_samples: return list(range(1, total_pages + 1))
        head_count = min(6, total_pages)
        tail_count = min(6, total_pages)
        head_pages = list(range(1, head_count + 1))
        tail_pages = list(range(max(total_pages - tail_count + 1, head_count + 1), total_pages + 1))
        mid_count = max_samples - len(head_pages) - len(tail_pages)
        if mid_count > 0:
            mid_start = head_count + 1
            mid_end = total_pages - tail_count
            if mid_end > mid_start:
                step = (mid_end - mid_start + 1) / (mid_count + 1)
                mid_pages = [int(mid_start + step * (i + 1)) for i in range(mid_count)]
                mid_pages = [p for p in mid_pages if p not in head_pages and p not in tail_pages]
            else: mid_pages = []
        else: mid_pages = []
        return sorted(set(head_pages + mid_pages + tail_pages))

    def _save_debug_image(self, image, pdf_path: str, page_number: int):
        if image is None: return
        try:
            pdf_name = Path(pdf_path).stem
            debug_dir = Path("/tmp/ocr_debug") / pdf_name
            debug_dir.mkdir(parents=True, exist_ok=True)
            image.save(debug_dir / f"page_{page_number:03d}.png")
        except: pass

    def extract_with_markers(self, pdf_path: str, prefix: str = "MAIN"):
        """
        메인 추출 로직
        """
        pages_text = []
        total_pages = 0
        ocr_count = 0
        
        # 통계 초기화
        self._gemini_ocr_used_pages = 0
        self._gemini_ocr_skipped_pages = 0

        # 1. 페이지 샘플링 계산
        sample_pages = None
        try:
            with pdfplumber.open(pdf_path) as pdf:
                total_pages = len(pdf.pages)
                if self.gemini_ocr_fallback:
                    sample_pages = self._calculate_sample_pages(total_pages, self.gemini_ocr_max_sample_pages)
                    _log(f"🎯 Gemini 샘플링: {len(sample_pages)}/{total_pages} 페이지", level="INFO")
        except Exception as e:
            _log(f"❌ PDF 열기 실패: {e}", level="ERROR")
            return {"full_text": "", "total_pages": 0, "gemini_fallback_used": False}

        # 2. 페이지별 순회
        with pdfplumber.open(pdf_path) as pdf:
            for page_idx, page in enumerate(pdf.pages, start=1):
                # A. 텍스트 레이어 추출 (가장 빠르고 정확, 0원)
                text = page.extract_text() or ""
                text_length = len(text.strip())

                # B. 텍스트가 부족하면 이미지 OCR 시도
                if text_length < self.min_text_length:
                    _log(f"page={page_idx} 텍스트 부족({text_length}자) -> 이미지 OCR 시도", level="DEBUG")
                    
                    # (1) ONNX OCR 시도 + 이미지 렌더링
                    ocr_text, pil_img = self._perform_ocr_on_page(pdf_path, page_idx)
                    
                    # 디버그 이미지 저장
                    self._save_debug_image(pil_img, pdf_path, page_idx)

                    if ocr_text and len(ocr_text) > 50:
                        text = ocr_text
                        ocr_count += 1
                        _log(f"✅ ONNX OCR 성공 ({len(text)}자)", level="INFO")
                    
                    # (2) ONNX 실패 시 Gemini Fallback
                    elif self.gemini_ocr_fallback and pil_img is not None:
                        if sample_pages and page_idx in sample_pages:
                            try:
                                buf = io.BytesIO()
                                pil_img.save(buf, format="PNG")
                                gem_text, usage = gemini_ocr_image_bytes(
                                    buf.getvalue(),
                                    language_hint="ko",
                                )
                                self._gemini_ocr_used_pages += 1
                                if gem_text and gem_text.strip():
                                    text = gem_text
                                    ocr_count += 1
                                    _log(f"✅ Gemini Vision 성공 ({len(text)}자)", level="INFO")
                                else:
                                    _log("⚠️ Gemini 결과 없음", level="WARNING")
                            except Exception as e:
                                _log(f"⚠️ Gemini 호출 실패: {e}", level="WARNING")
                        else:
                            self._gemini_ocr_skipped_pages += 1

                # 결과 저장
                title = text.split("\n")[0][:50] if text.strip() else f"Page {page_idx}"
                pages_text.append(f"[{prefix}-PAGE {page_idx}: {title}]")
                pages_text.append(text)
                pages_text.append("")

        if ocr_count:
            _log(f"✅ 총 OCR 처리 페이지: {ocr_count}", level="INFO")

        return {
            "full_text": "\n".join(pages_text),
            "total_pages": total_pages,
            "gemini_fallback_used": self._gemini_ocr_used_pages > 0,
        }

class ImageDescriptionGenerator:
    """통과된 이미지에 대한 상세 설명 생성 (2-4문장)"""
    
    
    def __init__(self):
        """이미지 설명 생성기 초기화"""
        self.total_tokens = 0  # ✅ 누적 토큰 수
        self.description_count = 0  # 생성한 설명 개수
        
        # ✅ Gemini 모델 초기화
        from .improved_hybrid_filter import get_global_model
        self.model = get_global_model()
        
        if self.model is None:
            print("      ⚠️  Warning: Gemini 모델 초기화 실패 - 이미지 설명 생성 불가", level="WARNING")
    def generate_description(
        self, 
        image_bytes: bytes, 
        adjacent_text: str,
        keywords: List[str],
        max_retries=3
    ) -> str:
        """
        Vision API로 이미지 상세 설명 생성
        재시도 로직 포함 (429 Rate Limit 대응)
        """
        import time
        
        for attempt in range(max_retries):
            try:
                mime_type = self._get_mime_type(image_bytes)
                image_part = Part.from_data(data=image_bytes, mime_type=mime_type)
                
                keyword_context = ', '.join(keywords[:10]) if keywords else "일반 학습 내용"
                
                prompt = f"""
이 이미지를 2-4문장으로 설명하세요.

강의 주제: {keyword_context}
주변 텍스트: "{adjacent_text}"

설명에 포함할 내용:
1. 이미지가 나타내는 주제/개념 (1문장)
2. 주요 구성 요소 2-3개 (1-2문장)
3. 핵심 정보나 패턴 (1문장)

제외할 내용:
- 세부 요소 전체 나열
- 불필요한 추측이나 해석

출력: 명확하고 간결한 2-4문장만.
"""
                # ✅ Vertex AI model 가져오기
                model = get_global_model()
                if model is None:
                    return "이미지 설명 생성 실패: Vertex AI model not initialized"

                response = model.generate_content([image_part, prompt])
                description = response.text.strip()
                
                # ✅ 토큰 사용량 추적 (usage_metadata 우선)
                tokens_added = 0
                try:
                    # Method 1: response.usage_metadata (가장 정확)
                    if hasattr(response, 'usage_metadata') and response.usage_metadata:
                        tokens_added = getattr(response.usage_metadata, 'total_token_count', 0)
                        if tokens_added > 0:
                            self.total_tokens += tokens_added
                            self.description_count += 1
                except Exception:
                    pass
                
                return description
                
            except Exception as e:
                error_msg = str(e)
                
                if "429" in error_msg or "Resource exhausted" in error_msg:
                    if attempt < max_retries - 1:
                        wait_time = (attempt + 1) * 3
                        _log(f"      ⚠️  Rate Limit, {wait_time}초 대기 중...", level="WARNING", end='', flush=True)
                        time.sleep(wait_time)
                        _log(" 재시도", level="WARNING")
                        continue
                    else:
                        return "이미지 설명 생성 실패: API rate limit exceeded"
                else:
                    return f"이미지 설명 생성 실패: {error_msg}"
        
        return "이미지 설명 생성 실패: Failed after all retries"
    
    def _get_mime_type(self, image_bytes: bytes) -> str:
        """이미지 바이너리에서 MIME 타입 감지"""
        if image_bytes.startswith(b'\xff\xd8'):
            return "image/jpeg"
        elif image_bytes.startswith(b'\x89PNG\r\n\x1a\n'):
            return "image/png"
        elif image_bytes.startswith(b'GIF87a') or image_bytes.startswith(b'GIF89a'):
            return "image/gif"
        elif image_bytes.startswith(b'RIFF') and image_bytes[8:12] == b'WEBP':
            return "image/webp"
        return "image/png"


class MetadataGenerator:
    """
    메타데이터 생성 노드
    
    주강의자료 + 보조자료 → metadata.json
    """
    
    def __init__(self):
        self.converter = None
        self.text_extractor = TextExtractor()
        self.image_filter = ImprovedHybridFilterPipeline(auto_extract_keywords=True)
        self.image_describer = ImageDescriptionGenerator()
        self.debug = True  # 🔧 DEBUG 항상 켜기 (원인 파악용)
            
    def _extract_page_title(self, slide_title: str, adjacent_text: str) -> str:
        """의미있는 페이지 제목 추출"""
        if slide_title and slide_title.strip() and slide_title.lower() != "no title":
            return slide_title.strip()[:50]
        
        if adjacent_text:
            lines = adjacent_text.strip().split('\n')
            for line in lines:
                line = line.strip()
                if len(line) > 3 and not line.startswith('☞'):
                    return line[:50]
        
        return "페이지 제목 없음"
    
    def generate(
        self,
        primary_file: str,
        supplementary_files: Optional[List[str]] = None,
        output_path: str = "output/metadata.json"
    ) -> str:
        """메타데이터 생성"""
        _log(f"\n{'='*120}")
        _log(f"🎯 메타데이터 생성 시작", level="INFO")
        _log(f"{'='*120}")
        _log(f"주강의자료: {primary_file}", level="INFO")
        if supplementary_files:
            _log(f"보조자료: {len(supplementary_files)}개")
            for i, supp in enumerate(supplementary_files, 1):
                _log(f"  {i}. {supp}")
        _log(f"{'='*120}\n")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            self.converter = DocumentConverterNode(output_dir=temp_dir)
            
            _log("📄 [1/3] 주강의자료 처리 중...", level="INFO")
            primary_metadata = self._process_primary_source(primary_file)
            
            _log("\n📚 [2/3] 보조자료 처리 중...", level="INFO")
            supplementary_metadata = []
            if supplementary_files:
                for i, supp_file in enumerate(supplementary_files[:3], 1):
                    try:
                        supp_meta = self._process_supplementary_source(supp_file, i)
                        supplementary_metadata.append(supp_meta)
                        _log(f"   ✅ 보조자료 {i} 처리 성공", level="INFO")
                    except Exception as e:
                        _log(f"   ⚠️ 보조자료 {i} 처리 실패 (계속 진행): {e}", level="WARNING", exc_info=True)
            else:
                _log("   ⚠️  보조자료 없음 (선택 사항)", level="INFO")
            
            _log("\n🔧 [3/3] 메타데이터 통합 중...", level="INFO")
            
            # ✅ Vision 토큰 통계 수집
            vision_tokens = {}
            if hasattr(self.image_filter, 'vision_tokens'):
                vision_tokens = self.image_filter.vision_tokens.copy()
                _log(f"   image_filter.vision_tokens = {vision_tokens}", level="DEBUG")
            
            # ✅ 이미지 설명 생성 토큰 추가
            _log(f"   image_describer.total_tokens = {self.image_describer.total_tokens}", level="DEBUG")
            _log(f"   image_describer.description_count = {self.image_describer.description_count}", level="DEBUG")
            
            if self.image_describer.total_tokens > 0:
                vision_tokens['image_description'] = self.image_describer.total_tokens
                vision_tokens['description_count'] = self.image_describer.description_count
                vision_tokens['total'] = vision_tokens.get('total', 0) + self.image_describer.total_tokens
                _log(f"   vision_tokens after adding image_description = {vision_tokens}", level="DEBUG")
            
            # ✅ 비용 계산
            if vision_tokens.get('total', 0) > 0:
                from .pricing import calculate_vision_cost, format_cost
                vision_cost = calculate_vision_cost(vision_tokens['total'])
                vision_tokens['cost_usd'] = vision_cost
            
            metadata = {
                "metadata_version": "1.0",
                "created_at": datetime.now().isoformat(),
                "primary_source": primary_metadata,
                "supplementary_sources": supplementary_metadata
            }
            
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, ensure_ascii=False, indent=2)
            
            _log(f"\n{'='*120}")
            _log(f"✅ 메타데이터 생성 완료!", level="INFO")
            _log(f"{'='*120}")
            _log(f"📁 출력 파일: {output_path}")
            _log(f"📊 주강의자료 페이지: {primary_metadata['total_pages']}개")
            _log(f"🖼️  필터링된 이미지: {len(primary_metadata['filtered_images'])}개")
            if supplementary_metadata:
                total_supp_pages = sum(s['total_pages'] for s in supplementary_metadata)
                _log(f"📚 보조자료 페이지: {total_supp_pages}개", level="INFO")
            
            # ✅ Vision 토큰 통계 출력
            if vision_tokens:
                _log(f"\n💰 Vision API 사용 통계:", level="INFO")
                if 'keyword_extraction' in vision_tokens:
                    _log(f"   📝 키워드 추출: {vision_tokens['keyword_extraction']:,} tokens", level="INFO")
                if 'image_filtering' in vision_tokens:
                    _log(f"   🔍 이미지 필터링: {vision_tokens['image_filtering']:,} tokens", level="INFO")
                if 'image_description' in vision_tokens:
                    _log(f"   📸 이미지 설명 생성: {vision_tokens['image_description']:,} tokens ({vision_tokens['description_count']}개)", level="INFO")
                if 'total' in vision_tokens:
                    _log(f"   📊 Total: {vision_tokens['total']:,} tokens", level="INFO")
                if 'cost_usd' in vision_tokens:
                    _log(f"   💵 비용: {format_cost(vision_tokens['cost_usd'])}", level="INFO")
            
            print(f"{'='*120}\n")
            
            # ✅ vision_tokens와 함께 반환
            return {
                "metadata_path": str(output_path),
                "vision_tokens": vision_tokens
            }
    
    def _process_primary_source(self, file_path: str) -> Dict[str, Any]:
        """
        주강의자료 처리
        ✅ TXT/URL 지원 추가
        ✅ PPTX 직접 텍스트 추출 (PDF 변환 없이)
        """
        file_path_str = str(file_path)
        
        # 원본 파일 타입 감지
        if file_path_str.startswith(('http://', 'https://')):
            original_file_type = 'url'
            file_path_obj = None
            display_name = file_path_str[:50]
        else:
            file_path_obj = Path(file_path)
            original_file_type = file_path_obj.suffix.lower().replace('.', '')
            display_name = file_path_obj.name
        
        _log(f"   📄 파일: {display_name} ({original_file_type})", level="INFO")
        
        # ✅ PPTX는 직접 텍스트 추출 (PDF 변환 시 한글 깨짐 방지)
        if original_file_type == 'pptx':
            _log(f"   📝 PPTX 직접 텍스트 추출 중... (PDF 변환 건너뜀)", level="INFO")
            from pptx import Presentation
            
            prs = Presentation(file_path_str)
            pages_text = []
            total_pages = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                total_pages += 1
                
                # 슬라이드 제목 추출
                title = "No Title"
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()[:50]
                
                # 페이지 마커
                pages_text.append(f"[MAIN-PAGE {slide_num}: {title}]")
                
                # 슬라이드 내용 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        pages_text.append(shape.text.strip())
                
                pages_text.append("")  # 슬라이드 구분
            
            full_text = "\n".join(pages_text)
            _log(f"   ✅ 텍스트 추출 완료: {len(full_text)}자, {total_pages}페이지", level="INFO")
            
            # 이미지는 PPTX 원본에서 추출
            _log(f"   🖼️  이미지 처리 중...", level="INFO")
            _log(f"      → PPTX 원본에서 직접 추출", level="INFO")
            self.image_filter.extract_keywords_from_document(file_path_str, text=full_text)
            keywords = self.image_filter.document_keywords
            all_images = self._extract_images_from_pptx(file_path_str)
            
        else:
            # 기존 방식: PDF 변환
            _log(f"   🔄 파일 처리 중...", level="INFO")
            processed_path = self.converter.convert(file_path_str)
            
            # 2. 텍스트 추출
            _log(f"   📝 텍스트 추출 중...", level="INFO")
            text_data = self.text_extractor.extract_with_markers(processed_path, prefix="MAIN")
            full_text = text_data['full_text']
            total_pages = text_data['total_pages']
            _log(f"   ✅ 텍스트 추출 완료: {len(full_text)}자", level="INFO")
            
            # 3. 이미지 필터링
            # 3. 이미지 필터링
            _log(f"   🖼️  이미지 처리 중...", level="INFO")
            
            # TXT/URL은 이미지 없음
            if original_file_type in ['txt', 'url']:
                _log(f"      → TXT/URL은 이미지 없음, 건너뛰기", level="INFO")
                all_images = []
                keywords = []
            
            elif original_file_type in ['docx', 'pdf']:
                _log(f"      → PDF에서 이미지 추출", level="INFO")
                self.image_filter.extract_keywords_from_document(processed_path, text=full_text)
                keywords = self.image_filter.document_keywords
                extractor = UniversalImageExtractor()
                
                # ✅ Gemini Fallback 사용 여부 전달
                gemini_used = text_data.get('gemini_fallback_used', False)
                all_images = extractor.extract(processed_path, skip_ocr=gemini_used)
            
            else:
                _log(f"   ⚠️  지원하지 않는 형식: {original_file_type}", level="WARNING")
                all_images = []
                keywords = []
        
        # 4. 필터링 실행 (공통)
        filtered_images = []
        if all_images:
            _log(f"   🔍 {len(all_images)}개 이미지 발견, 필터링 시작...")

            for img_meta in all_images:
                decision, reason = self.image_filter.step1_rule_check(img_meta)
                
                if decision == "INCLUDE":
                    # ✅ V3: Rule 통과도 AI로 검증 + 설명 생성
                    result = self.image_filter.unified_vision_check(img_meta)
                    
                    if result["is_core"]:
                        img_meta.is_core_content = True
                        img_meta.description = result["description"] or ""
                        img_meta.filter_reason = f"Rule+AI: {result['reason']}"
                        filtered_images.append(img_meta)
                    
                elif decision == "PENDING":
                    # ✅ V3: unified_vision_check 사용 (필터링 + 설명 통합)
                    result = self.image_filter.unified_vision_check(img_meta)
                    
                    if result["is_core"]:
                        img_meta.is_core_content = True
                        img_meta.description = result["description"] or ""
                        img_meta.filter_reason = result["reason"]
                        filtered_images.append(img_meta)
            
            _log(f"   ✅ 필터링 완료: {len(filtered_images)}개 선택")
        
        # 5. 이미지 메타데이터 구성
        filtered_image_metadata = []
        
        if filtered_images:
            # ✅ V3: 설명이 이미 포함되어 있음 (unified_vision_check에서 생성)
            _log(f"   ✅ 이미지 메타데이터 구성 중... ({len(filtered_images)}개)", level="INFO")
            
            for i, img_meta in enumerate(filtered_images, 1):
                page_title = self._extract_page_title(
                    img_meta.slide_title,
                    img_meta.adjacent_text
                )
                
                # ✅ description은 이미 img_meta.description에 존재!
                filtered_image_metadata.append({
                    "image_id": img_meta.image_id.replace("S", "MAIN_P").replace("P", "MAIN_P"),
                    "page_number": img_meta.slide_number,
                    "page_title": page_title,
                    "description": img_meta.description or "설명 없음",  # ✅ 이미 생성됨
                    "filter_stage": "1차 (Rule+AI)" if "Rule+AI" in img_meta.filter_reason else "2차 (AI)",
                    "area_percentage": img_meta.area_percentage
                })
            
            _log(f"   ✅ 메타데이터 구성 완료: {len(filtered_image_metadata)}개", level="INFO")
            _log(f"   ⚡ 최적화: 통합 Vision API로 설명 생성 중복 제거", level="INFO")

        # 6. 통계
        total_images = len(all_images)
        passed_images = len(filtered_images)
        
        return {
            "role": "main",
            "filename": display_name if original_file_type == 'url' else (file_path_obj.name if file_path_obj else display_name),
            "file_type": original_file_type,
            "total_pages": total_pages,
            "content": {
                "full_text": full_text
            },
            "filtered_images": filtered_image_metadata,
            "statistics": {
                "total_images_found": total_images,
                "images_passed": passed_images,
                "filter_rate": passed_images / total_images if total_images > 0 else 0
            }
        }
    
    def _process_supplementary_source(self, file_path: str, order: int) -> Dict[str, Any]:
        """
        보조자료 처리
        ✅ PPTX 직접 텍스트 추출 (PDF 변환 없이)
        """
        file_path_str = str(file_path)
        
        # URL과 파일 구분
        if file_path_str.startswith(('http://', 'https://')):
            file_type = 'url'
            display_name = 'Web Content'
            file_path_obj = None
        else:
            file_path_obj = Path(file_path)
            file_type = file_path_obj.suffix.lower().replace('.', '')
            display_name = file_path_obj.name
        
        _log(f"   📚 보조자료 {order}: {display_name} ({file_type})")
        
        # ✅ PPTX는 직접 텍스트 추출 (PDF 변환 건너뜀)
        if file_type == 'pptx':
            print(f"      📝 PPTX 직접 텍스트 추출 중... (PDF 변환 건너뜀)")
            from pptx import Presentation
            
            prs = Presentation(file_path_str)
            pages_text = []
            total_pages = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                total_pages += 1
                
                # 슬라이드 제목 추출
                title = "No Title"
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title = slide.shapes.title.text.strip()[:50]
                
                # 페이지 마커
                pages_text.append(f"[SUPP{order}-PAGE {slide_num}: {title}]")
                
                # 슬라이드 내용 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        pages_text.append(shape.text.strip())
                
                pages_text.append("")  # 슬라이드 구분
            
            full_text = "\n".join(pages_text)
            print(f"      ✅ 완료 ({total_pages}페이지)")
            
        else:
            # 기존 방식: PDF 변환
            print(f"      🔄 PDF 변환 중...")
            pdf_path = self.converter.convert(file_path_str)
            
            print(f"      📝 텍스트 추출 중...")
            text_data = self.text_extractor.extract_with_markers(pdf_path, prefix=f"SUPP{order}")
            
            full_text = text_data['full_text']
            total_pages = text_data['total_pages']
            
            print(f"      ✅ 완료 ({total_pages}페이지)")
        
        return {
            "order": order,
            "filename": display_name,
            "file_type": file_type,
            "total_pages": total_pages,
            "content": {
                "full_text": full_text
            }
        }
    
    def _extract_images_from_pptx(self, pptx_path: str) -> List[ImageMetadata]:
        """PPTX에서 이미지 메타데이터 추출"""
        extractor = UniversalImageExtractor()
        return extractor.extract(pptx_path)


# CLI 인터페이스
if __name__ == "__main__":
    import sys
    
    _log("\n" + "="*120)
    _log("🎯 Metadata Generator Node (V2 - pdfplumber)")
    _log("="*120)
    
    if len(sys.argv) < 2:
        _log("\n사용법:")
        _log("  python metadata_generator_node.py <주강의자료> [보조1] [보조2] [보조3]")
        _log("\n예시:")
        _log("  python metadata_generator_node.py 중등국어1.pptx")
        _log("  python metadata_generator_node.py notes.txt")
        _log("  python metadata_generator_node.py https://example.com/article")
        _log("\n✅ 지원 형식: PPTX, DOCX, PDF, TXT, URL")
        _log("="*120 + "\n")
        sys.exit(1)
    
    primary_file = sys.argv[1]
    supplementary_files = sys.argv[2:5] if len(sys.argv) > 2 else None
    
    if not primary_file.startswith('http') and not os.path.exists(primary_file):
        _log(f"\n❌ 주강의자료를 찾을 수 없습니다: {primary_file}")
        sys.exit(1)
    
    if supplementary_files:
        for supp in supplementary_files:
            if not supp.startswith('http') and not os.path.exists(supp):
                _log(f"\n❌ 보조자료를 찾을 수 없습니다: {supp}")
                sys.exit(1)
    
    try:
        generator = MetadataGenerator()
        output_path = generator.generate(
            primary_file=primary_file,
            supplementary_files=supplementary_files,
            output_path="output/metadata.json"
        )
        
        _log(f"✅ 성공!")
        _log(f"📁 {output_path}")
        
    except Exception as e:
        _log(f"\n❌ 에러 발생: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
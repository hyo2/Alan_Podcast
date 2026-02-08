"""
Improved Hybrid Filter V4
==========================

V4 변경사항:
- V2의 유연한 인증 로직 추가 (환경 변수 기반)
- V3의 pdfplumber (MIT) 유지
- 색상 복잡도 필터 유지
- 인증 실패 시 graceful degradation

핵심 기능:
- PyMuPDF (AGPL) → pdfplumber (MIT) 전환
- 라이선스 문제 해결
- 색상 복잡도 필터 (텍스트 상자 배경 제거)
- 환경 변수 기반 인증 (프로덕션 대응)
"""

import os
import textwrap
import json
from dataclasses import dataclass
from typing import List, Dict
from pptx import Presentation
from vertexai.generative_models import Part
import logging
import sys
logger = logging.getLogger(__name__)

def _log(*args, level: str | None = None, exc_info: bool = False, end: str = '\n', flush: bool = False) -> None:
    """
     logger 기반 로그 (환경별 LOG_LEVEL 적용).
     - 기본 level: DEBUG
     - end/flush를 쓰는 진행형 출력(end != '\\n' 또는 flush=True)은 print 유지
    """
    msg = " ".join(str(a) for a in args).rstrip() if args else ""
    # 진행형 출력은 그대로 stdout로 (기존 UX 유지)
    if end != "\n" or flush:
        print(msg, end=end, flush=flush)
        return

    lvl = (level or "DEBUG").upper()
    if lvl == "DEBUG":
        logger.debug(msg, exc_info=exc_info)
    elif lvl == "INFO":
        logger.info(msg, exc_info=exc_info)
    elif lvl in ("WARN", "WARNING"):
        logger.warning(msg, exc_info=exc_info)
    elif lvl == "ERROR":
        logger.error(msg, exc_info=exc_info)
    elif lvl in ("CRITICAL", "FATAL"):
        logger.critical(msg, exc_info=exc_info)
    else:
        logger.debug(msg, exc_info=exc_info)

def _resolve_vertex_sa_file() -> str | None:
    # 프로젝트에서 쓰는 키 우선순위
    # NOTE: VERTEX_AI_SERVICE_ACCOUNT_JSON(=JSON 문자열)은 main.py의 patch_vertex_ai_env()에서
    # 파일로 변환 후 GOOGLE_APPLICATION_CREDENTIALS로 연결되므로 여기서는 "경로"만 확인한다.
    for key in ("VERTEX_AI_SERVICE_ACCOUNT_FILE", "GOOGLE_APPLICATION_CREDENTIALS"):
        p = os.getenv(key)
        if p and os.path.exists(p):
            return p
    return None

def get_vertex_text_model():
    """
    키워드 추출/이미지 판단(vision)에서 쓰는 Gemini 모델 lazy init.
    - 인증 파일 없으면 None 반환 (로컬 데모에서 vision만 스킵 가능)
    """
    try:
        sa_file = _resolve_vertex_sa_file()
        if not sa_file:
            logger.warning("ℹ️ Vertex 서비스 계정 파일이 없어 Gemini 호출을 스킵합니다.")
            return None

        os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = sa_file

        import vertexai
        from vertexai.generative_models import GenerativeModel

        project_id = os.getenv("VERTEX_AI_PROJECT_ID")
        location = os.getenv("VERTEX_AI_REGION", "us-central1")

        if project_id:
            vertexai.init(project=project_id, location=location)
        else:
            vertexai.init(location=location)

        model_name = os.getenv("VERTEX_AI_MODEL_TEXT", "gemini-2.5-flash")
        return GenerativeModel(model_name)

    except Exception as e:
        logger.exception(f"Vertex/Gemini 초기화 실패: {e}")
        return None
    
model = None

def get_global_model():
    global model
    if model is None:
        model = get_vertex_text_model()
    return model

@dataclass
class ImageMetadata:
    image_id: str
    slide_number: int
    area_percentage: float
    left: float
    top: float
    adjacent_text: str
    slide_title: str
    image_bytes: bytes = None
    is_core_content: bool = False
    filter_reason: str = ""

# 1. 통합 이미지 추출기 (PPTX + PDF 지원)
class UniversalImageExtractor:
    """
    모든 형식에서 이미지 메타데이터 추출
    V3: pdfplumber (MIT) 사용
    """
    
    def extract(self, file_path: str) -> List[ImageMetadata]:
        from pathlib import Path
        
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx':
            return self._extract_from_pptx(file_path)
        elif ext == '.pdf':
            return self._extract_from_pdf_v3(file_path)  # ✅ v3로 변경
        else:
            raise ValueError(f"지원하지 않는 형식: {ext}")
    
    def _extract_from_pptx(self, pptx_path: str) -> List[ImageMetadata]:
        """PPTX에서 이미지 추출 (기존 방식)"""
        if not os.path.exists(pptx_path):
            return []
        
        prs = Presentation(pptx_path)
        metadata_list = []
        slide_width, slide_height = prs.slide_width.inches, prs.slide_height.inches
        slide_area = slide_width * slide_height

        for s_idx, slide in enumerate(prs.slides, 1):
            slide_title = slide.shapes.title.text if slide.shapes.title else "No Title"
            all_text = " ".join([s.text for s in slide.shapes if hasattr(s, "text")])
            
            img_idx = 1
            for shape in slide.shapes:
                if shape.shape_type == 13 or hasattr(shape, 'image'):
                    w, h = shape.width.inches, shape.height.inches
                    area_pct = ((w * h) / slide_area) * 100
                    metadata_list.append(ImageMetadata(
                        image_id=f"S{s_idx:02d}_IMG{img_idx:03d}",
                        slide_number=s_idx,
                        area_percentage=area_pct,
                        left=shape.left.inches,
                        top=shape.top.inches,
                        adjacent_text=all_text.replace('\n', ' ').strip(),
                        slide_title=slide_title,
                        image_bytes=shape.image.blob
                    ))
                    img_idx += 1
        
        return metadata_list
    
    def _safe_parse_paddleocr_result(self, ocr_result):
        """
        PaddleOCR v3~v5 결과를 안전하게 파싱
        반환: List[Dict] -> {text, bbox, confidence}
        """
        parsed = []

        if not ocr_result:
            return parsed

        # PaddleOCR는 보통 [result] 형태로 한 번 더 감싸짐
        if isinstance(ocr_result, list) and len(ocr_result) == 1 and isinstance(ocr_result[0], list):
            ocr_result = ocr_result[0]

        for item in ocr_result:
            try:
                # ✅ Case 1: dict 형태 (PaddleOCR v5 / PaddleX)
                if isinstance(item, dict):
                    text = item.get("text", "").strip()
                    bbox = item.get("points") or item.get("bbox")
                    conf = item.get("confidence") or item.get("score")

                # ✅ Case 2: classic list 형태 [[bbox], (text, score)]
                elif isinstance(item, (list, tuple)):
                    # ("text", score) 형태
                    if len(item) == 2 and isinstance(item[0], str):
                        text = item[0].strip()
                        bbox = None
                        conf = item[1]

                    # [[x,y]... , ("text", score)]
                    elif len(item) >= 2 and isinstance(item[1], (list, tuple)):
                        bbox = item[0]
                        text = item[1][0].strip() if len(item[1]) > 0 else ""
                        conf = item[1][1] if len(item[1]) > 1 else None
                    else:
                        continue
                else:
                    continue

                if not text:
                    continue

                parsed.append({
                    "text": text,
                    "bbox": bbox,
                    "confidence": conf
                })

            except Exception:
                # 파싱 실패해도 전체 OCR은 살리기
                continue

        return parsed

    def _normalize_ocr_image(self, pil_img):
        """
        PaddleOCR 안정화를 위한 이미지 정규화
        - RGBA → RGB
        - Grayscale → RGB
        - numpy contiguous 보장
        """
        from PIL import Image
        import numpy as np

        if pil_img.mode != "RGB":
            pil_img = pil_img.convert("RGB")

        img_array = np.array(pil_img)

        # numpy contiguous 보장 (중요)
        if not img_array.flags['C_CONTIGUOUS']:
            img_array = np.ascontiguousarray(img_array)

        return img_array


    def _extract_text_with_ocr(self, pdf_path: str, page_num: int, min_length: int = 100) -> str:
        """
        페이지에서 텍스트 추출 (필요시 OCR)
        V3: pdfplumber + pypdfium2 + PaddleOCR
        """
        # ===== 1. pdfplumber로 먼저 시도 =====
        try:
            import pdfplumber
            with pdfplumber.open(pdf_path) as pdf:
                page = pdf.pages[page_num]
                text = page.extract_text() or ""
                text_length = len(text.strip())
                
                if text_length >= min_length:
                    return text
        except:
            text_length = 0
        
        # ===== 2. 텍스트 부족 → OCR 실행 =====
        try:
            from paddleocr import PaddleOCR
            import numpy as np
            from pypdfium2 import PdfDocument

            if not hasattr(self, '_ocr_engine'):
                os.environ['FLAGS_log_level'] = '3'
                os.environ['PPOCR_SHOW_LOG'] = 'False'

                _log(f"      → PaddleOCR 초기화 중...")
                self._ocr_engine = PaddleOCR(lang='korean', use_angle_cls=True)

            # ✅ pypdfium2로 해당 페이지만 렌더링 (0-indexed page_num)
            pdf = PdfDocument(pdf_path)
            if page_num < 0 or page_num >= len(pdf):
                return text

            page = pdf[page_num]
            target_dpi = 150
            scale = target_dpi / 72.0
            bitmap = page.render(scale=scale)
            pil_img = bitmap.to_pil()

            img_array = self._normalize_ocr_image(pil_img)

            # ===== OCR 실행 (안정화 버전) =====
            result = None

            # 페이지 OCR에서는 cls=True 금지
            try:
                result = self._ocr_engine.ocr(img_array)
            except Exception as e1:
                _log(f"      ⚠️ ocr() 실패, predict() 시도: {e1}")
                try:
                    if hasattr(self._ocr_engine, "predict"):
                        result = self._ocr_engine.predict(img_array)
                except Exception as e2:
                    _log(f"      ❌ OCR 완전 실패: {e2}")
                    return text

            parsed = self._safe_parse_paddleocr_result(result)

            if parsed:
                lines = [p["text"] for p in parsed]
                ocr_result = "\n".join(lines)
                _log(f"      → OCR 완료: {text_length}자 → {len(ocr_result)}자")
                return ocr_result if ocr_result else text


        except ImportError as e:
            _log(f"      → PaddleOCR/pypdfium2 미설치, 텍스트만 사용: {e}")
            return text
        except Exception as e:
            _log(f"      ⚠️  OCR 실패: {e}")
            return text

        return text

    
    def _extract_text_bboxes_with_ocr(self, pdf_path: str, page_num: int) -> List[Dict]:
        """
        페이지에서 텍스트 bbox 추출 (OCR 활용)
        
        Returns:
            [{'x0', 'top', 'x1', 'bottom'}, ...]
        """
        text_bboxes = []
        
        # ===== 1. pdfplumber로 먼저 시도 =====
        try:
            import pdfplumber
            with pdfplumber.open(pdf_path) as pdf:
                page = pdf.pages[page_num]
                chars = page.chars
                
                if chars and len(chars) > 0:
                    # 텍스트 레이어가 있음
                    for char in chars:
                        text_bboxes.append({
                            'x0': char['x0'],
                            'top': char['top'],
                            'x1': char['x1'],
                            'bottom': char['bottom']
                        })
                    
                    _log(f"      → pdfplumber로 {len(text_bboxes)}개 문자 bbox 추출")
                    return text_bboxes
        except:
            pass
        
        # ===== 2. 텍스트 레이어 없음 → OCR로 bbox 추출 =====
        try:
            from paddleocr import PaddleOCR
            import numpy as np
            from pypdfium2 import PdfDocument

            if not hasattr(self, '_ocr_engine'):
                os.environ['FLAGS_log_level'] = '3'
                os.environ['PPOCR_SHOW_LOG'] = 'False'
                self._ocr_engine = PaddleOCR(lang='korean', use_angle_cls=True)

            # ✅ pypdfium2로 해당 페이지 렌더링
            pdf = PdfDocument(pdf_path)
            if page_num < 0 or page_num >= len(pdf):
                return []

            page = pdf[page_num]
            target_dpi = 150
            scale = target_dpi / 72.0
            bitmap = page.render(scale=scale)
            pil_img = bitmap.to_pil()

            img_array = self._normalize_ocr_image(pil_img)

            # ===== OCR 실행 (버전별 대응) =====
            result = None

            try:
                result = self._ocr_engine.ocr(img_array)
            except Exception as e1:
                _log(f"      ⚠️ bbox OCR ocr() 실패, predict() 시도: {e1}")
                try:
                    if hasattr(self._ocr_engine, "predict"):
                        result = self._ocr_engine.predict(img_array)
                except Exception as e2:
                    _log(f"      ❌ bbox OCR 완전 실패: {e2}")
                    return []

            parsed = self._safe_parse_paddleocr_result(result)

            for item in parsed:
                bbox_points = item["bbox"]
                if not bbox_points:
                    continue

                try:
                    x_coords = [p[0] for p in bbox_points if len(p) == 2]
                    y_coords = [p[1] for p in bbox_points if len(p) == 2]

                    if not x_coords or not y_coords:
                        continue

                    text_bboxes.append({
                        'x0': min(x_coords),
                        'top': min(y_coords),
                        'x1': max(x_coords),
                        'bottom': max(y_coords)
                    })
                except Exception:
                    continue

            if text_bboxes:
                _log(f"      → OCR로 {len(text_bboxes)}개 텍스트 bbox 추출")

            return text_bboxes

        except ImportError as e:
            _log(f"      ⚠️  OCR bbox 추출 불가(PaddleOCR/pypdfium2 미설치): {e}")
            return []
        except Exception as e:
            _log(f"      ⚠️  텍스트 bbox 추출 실패: {e}")
            return []

    
    def _calculate_text_overlap(self, img_bbox: tuple, text_bboxes: List[Dict]) -> float:
        """
        이미지와 텍스트의 중첩 비율 계산
        
        Args:
            img_bbox: (x0, top, x1, bottom)
            text_bboxes: [{'x0', 'top', 'x1', 'bottom'}, ...]
        
        Returns:
            중첩 비율 (0.0 ~ 1.0)
        """
        if not text_bboxes:
            return 0.0
        
        img_x0, img_top, img_x1, img_bottom = img_bbox
        img_area = (img_x1 - img_x0) * (img_bottom - img_top)
        
        if img_area <= 0:
            return 0.0
        
        overlap_area = 0.0
        
        for text_bbox in text_bboxes:
            # 교집합 계산
            x0 = max(img_x0, text_bbox['x0'])
            top = max(img_top, text_bbox['top'])
            x1 = min(img_x1, text_bbox['x1'])
            bottom = min(img_bottom, text_bbox['bottom'])
            
            if x0 < x1 and top < bottom:
                overlap_area += (x1 - x0) * (bottom - top)
        
        overlap_ratio = overlap_area / img_area
        
        return overlap_ratio
    
    def _calculate_color_complexity(self, image_bytes) -> int:
        """
        이미지의 색상 복잡도 계산 (고유 색상 수)
        
        텍스트 상자 배경: 10-300개 (단조로운 색상)
        진짜 콘텐츠: 500+ 개 (복잡한 색상)
        
        Args:
            image_bytes: 이미지 바이너리 데이터
        
        Returns:
            고유 색상 수 (0 ~ 10000+)
        """
        try:
            from PIL import Image
            import io
            
            # 바이너리 → PIL Image
            img = Image.open(io.BytesIO(image_bytes))
            
            # RGB 변환
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # 너무 크면 리사이즈 (속도 향상)
            max_size = 500
            if img.width > max_size or img.height > max_size:
                ratio = min(max_size / img.width, max_size / img.height)
                new_size = (int(img.width * ratio), int(img.height * ratio))
                img = img.resize(new_size, Image.Resampling.LANCZOS)
            
            # 고유 색상 수 계산
            colors = img.getcolors(maxcolors=10000)
            
            if colors:
                unique_colors = len(colors)
            else:
                # 10000개 이상 색상
                unique_colors = 10000
            
            return unique_colors
        
        except Exception as e:
            logger.warning(f"색상 분석 실패: {e}")
            return 10000  # 실패 시 복잡한 이미지로 간주
    
    def _extract_page_title(self, page_text: str) -> str:
        """페이지 제목 추출"""
        lines = page_text.strip().split('\n')
        for line in lines:
            line = line.strip()
            if len(line) > 3 and not line.startswith('☞'):
                return line[:50]
        return "페이지 제목 없음"
    
    def _extract_from_pdf_v3(self, pdf_path: str) -> List[ImageMetadata]:
        """
        PDF에서 이미지 추출 (V3: pdfplumber 사용)
        
        핵심 변경:
        - PyMuPDF → pdfplumber (MIT 라이선스)
        - 기능 동일하게 유지
        """
        try:
            import pdfplumber  # ✅ pdfplumber 사용
        except ImportError:
            _log("   ❌ pdfplumber가 설치되지 않았습니다.")
            _log("   pip install pdfplumber")
            return []
        
        if not os.path.exists(pdf_path):
            return []
        
        metadata_list = []
        
        # 필터링 기준
        MIN_WIDTH = 40
        MIN_HEIGHT = 40
        MIN_AREA_PCT = 3.0      # 3% 미만: 레이블/아이콘
        MAX_AREA_PCT = 90.0     # 90% 이상: 배경
        MIN_PIXEL_AREA = 1000
        MAX_ASPECT_RATIO = 6.0  # 6:1 이상: 제목/텍스트
        
        total_images = 0
        filtered_background = 0
        filtered_aspect = 0
        filtered_area = 0
        filtered_size = 0
        filtered_text_overlap = 0  # ✅ 추가
        
        try:
            # ===== pdfplumber로 PDF 열기 =====
            with pdfplumber.open(pdf_path) as pdf:
                
                for page_num, page in enumerate(pdf.pages):
                    # 페이지 정보
                    page_width = page.width
                    page_height = page.height
                    page_area = page_width * page_height
                    
                    # 텍스트 추출 (OCR 포함)
                    page_text = self._extract_text_with_ocr(pdf_path, page_num, min_length=100)
                    page_title = self._extract_page_title(page_text)
                    
                    # ===== 텍스트 bbox 추출 (중첩 체크용) =====
                    text_bboxes = self._extract_text_bboxes_with_ocr(pdf_path, page_num)
                    
                    # ===== pdfplumber로 이미지 목록 가져오기 =====
                    images = page.images
                    total_images += len(images)
                    
                    _log(f"      [P{page_num+1}] 총 {len(images)}개 이미지 발견")
                    
                    for img in images:
                        try:
                            # ===== bbox 정보 (pdfplumber 형식) =====
                            x0 = img['x0']
                            top = img['top']
                            x1 = img['x1']
                            bottom = img['bottom']
                            
                            width = x1 - x0
                            height = bottom - top
                            area_pct = (width * height) / page_area * 100
                            
                            debug_msg = f"      [P{page_num+1}] {area_pct:.1f}%"
                            
                            # ===== 필터 1: 배경 제외 (90% 이상) =====
                            if area_pct > MAX_AREA_PCT:
                                filtered_background += 1
                                _log(debug_msg + f" → 배경 제외 ❌")
                                continue
                            
                            # ===== 필터 2: 가로세로비 =====
                            if width > 0 and height > 0:
                                aspect_ratio = max(width, height) / min(width, height)
                                if aspect_ratio > MAX_ASPECT_RATIO:
                                    filtered_aspect += 1
                                    _log(debug_msg + f" → 가로세로비 제외 ({aspect_ratio:.1f}:1) ❌")
                                    continue
                            
                            # ===== 필터 3: 작은 면적 =====
                            pixel_area = width * height
                            if pixel_area < MIN_PIXEL_AREA:
                                filtered_area += 1
                                _log(debug_msg + f" → 작은 면적 제외 ❌")
                                continue
                            
                            # ===== 필터 4: 절대 크기 =====
                            if width < MIN_WIDTH or height < MIN_HEIGHT:
                                filtered_size += 1
                                _log(debug_msg + f" → 작은 크기 제외 ❌")
                                continue
                            
                            # ===== 필터 5: 상대 크기 =====
                            if area_pct < MIN_AREA_PCT:
                                filtered_size += 1
                                _log(debug_msg + f" → 상대 크기 제외 ({area_pct:.1f}%) ❌")
                                continue
                            
                            # ===== 통과! =====
                            _log(debug_msg + " → 최종 추출 ✅✅✅")
                            
                            # ===== 필터 6: 텍스트 중첩 + 색상 복잡도 체크 ⭐⭐⭐ =====
                            # 이미지 바이너리 추출
                            stream = img.get('stream')
                            
                            if stream:
                                if hasattr(stream, 'get_data'):
                                    image_bytes = stream.get_data()
                                elif hasattr(stream, 'rawdata'):
                                    image_bytes = stream.rawdata
                                else:
                                    _log(debug_msg + " → 바이너리 추출 실패 ⚠️")
                                    continue
                            else:
                                _log(debug_msg + " → stream 없음 ⚠️")
                                continue
                            
                            # 텍스트 중첩 계산
                            img_bbox = (x0, top, x1, bottom)
                            overlap_ratio = self._calculate_text_overlap(img_bbox, text_bboxes)
                            
                            # 색상 복잡도 계산
                            color_count = self._calculate_color_complexity(image_bytes)
                            
                            # 판단 로직 (색상 + 중첩)
                            is_textbox = False
                            filter_reason = ""
                            
                            # 규칙 1: 단조로운 색상 (< 300개) → 텍스트 상자 가능성
                            if color_count < 300:
                                if overlap_ratio >= 0.05:  # 5% 이상 중첩
                                    is_textbox = True
                                    filter_reason = f"단조색상({color_count}개)+중첩({overlap_ratio*100:.0f}%)"
                                elif area_pct >= 15.0:
                                    is_textbox = True
                                    filter_reason = f"단조색상({color_count}개)+대형"
                            
                            # 규칙 2: 복잡한 색상 (>= 500개) → 진짜 콘텐츠 가능성
                            elif color_count >= 500:
                                if overlap_ratio >= 0.30:  # 30% 이상만 제외
                                    is_textbox = True
                                    filter_reason = f"고중첩({overlap_ratio*100:.0f}%)"
                                # else: 통과
                            
                            # 규칙 3: 중간 복잡도 (300-500개) → 중첩 비율로 판단
                            else:
                                if overlap_ratio >= 0.15:  # 15% 이상
                                    is_textbox = True
                                    filter_reason = f"중간색상({color_count}개)+중첩({overlap_ratio*100:.0f}%)"
                            
                            # 결과 처리
                            if is_textbox:
                                filtered_text_overlap += 1
                                _log(debug_msg + f" → 텍스트상자 제외 ({filter_reason}) ❌")
                                continue
                            
                            # 최종 통과 - 메타데이터 저장
                            
                            metadata_list.append(ImageMetadata(
                                image_id=f"P{page_num+1:02d}_IMG{len(metadata_list)+1:03d}",
                                slide_number=page_num + 1,
                                area_percentage=area_pct,
                                left=x0,
                                top=top,
                                adjacent_text=page_text.replace('\n', ' ').strip(),
                                slide_title=page_title,
                                image_bytes=image_bytes
                            ))
                        
                        except Exception as e:
                            _log(f"      ⚠️ 이미지 처리 실패: {e}")
                            continue
        
        except Exception as e:
            _log(f"   ❌ PDF 처리 실패: {e}")
            import traceback
            traceback.print_exc()
            return []
        
        # 통계
        _log(f"\n   📊 PDF 이미지 분석:")
        _log(f"      - 전체 이미지: {total_images}개")
        _log(f"   🔍 필터링 통계:")
        _log(f"      - 배경 제외: {filtered_background}개")
        _log(f"      - 가로세로비: {filtered_aspect}개")
        _log(f"      - 작은 면적: {filtered_area}개")
        _log(f"      - 작은 크기: {filtered_size}개")
        _log(f"      - 텍스트 상자 (색상+중첩): {filtered_text_overlap}개")  # ✅ 추가
        _log(f"   ✅ 최종 추출: {len(metadata_list)}개 이미지\n")
        
        return metadata_list


# 2. 개선된 하이브리드 필터 파이프라인
class ImprovedHybridFilterPipeline:
    def __init__(self, auto_extract_keywords: bool = True):
        self.auto_extract = auto_extract_keywords
        
        self.UNIVERSAL_PATTERNS = [
            '학습', '활동', '문제', '예제', '연습',
            '생각', '알아보', '살펴보', '정리',
            '목표', '개념', '원리', '법칙', '정의',
            '단원', '차시',
            '그림', '도표', '표', '차트', '그래프',
            '예시', '사례', '모형', '구조'
        ]
        
        self.DECORATION_PATTERNS = [
            '로고', 'logo', '출처', '참고', '아이콘', 'icon'
        ]
        
        self.document_keywords = []
        
        self.model = get_global_model()

    def extract_keywords_from_document(self, file_path: str):
        """문서에서 자동으로 키워드 추출"""
        if not self.auto_extract:
            return
        
        from pathlib import Path
        
        _log("📚 문서 분석하여 키워드 자동 추출 중...")
        
        ext = Path(file_path).suffix.lower()
        all_text = []
        
        if ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
        
        elif ext == '.pdf':
            import pdfplumber  # ✅ pdfplumber 사용
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            all_text.append(text)
            except Exception as e:
                _log(f"   ⚠️ PDF 텍스트 추출 실패, 범용 패턴만 사용")
                return
        
        else:
            _log(f"   ⚠️ 지원하지 않는 형식: {ext}")
            return
        
        full_text = "\n".join(all_text)[:5000]
        
        prompt = f"""
다음 강의 자료에서 **핵심 키워드 20개**를 추출하세요.

# 문서 내용
{full_text}

# 조건
- 개념어, 전문 용어, 주제어만 포함
- JSON 형식: {{"keywords": ["키워드1", "키워드2", ...]}}
"""
        
        if self.model is None:
            _log("   ⚠️ Gemini 모델 초기화 실패(인증 없음). 키워드 자동 추출 스킵.")
            self.document_keywords = []
            return

        try:
            response = self.model.generate_content(prompt)
            text = response.text.strip()
            
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0].strip()
            elif "```" in text:
                text = text.split("```")[1].split("```")[0].strip()
            
            data = json.loads(text)
            self.document_keywords = data.get("keywords", [])
            
            _log(f"   ✅ 추출된 키워드: {', '.join(self.document_keywords[:10])}")
        
        except Exception as e:
            _log(f"   ⚠️ 자동 추출 실패, 범용 패턴만 사용")
            self.document_keywords = []

    def step1_rule_check(self, meta: ImageMetadata):
        """규칙 기반 1차 필터"""
        context = f"{meta.slide_title} {meta.adjacent_text}".lower()
        
        has_deco = any(kw in context for kw in self.DECORATION_PATTERNS)
        is_corner = (meta.left < 1.0 and meta.top < 1.0) or (meta.left > 8.0 and meta.top < 1.0)
        
        if is_corner and meta.area_percentage < 5.0 and not any(kw in context for kw in self.UNIVERSAL_PATTERNS):
            return "EXCLUDE", "Static Decoration (Corner)"
        
        if has_deco and meta.area_percentage < 8.0:
            return "EXCLUDE", "Decorative element"
        
        has_universal = any(p in context for p in self.UNIVERSAL_PATTERNS)
        has_document_kw = any(kw in context for kw in self.document_keywords)
        
        if meta.area_percentage > 15.0 and (has_universal or has_document_kw):
            return "INCLUDE", f"Core content ({meta.area_percentage:.1f}% + pattern)"
        
        if has_document_kw and meta.area_percentage > 10.0:
            matched = [kw for kw in self.document_keywords if kw in context]
            return "INCLUDE", f"Document keyword: {', '.join(matched[:2])}"
        
        return "PENDING", "Requires AI Vision Check"

    def step2_gemini_check(self, meta: ImageMetadata, max_retries=3):
        """AI Vision으로 2차 판단"""
        import time
        
        
        if self.model is None:
            return "DISCARD: Gemini unavailable (no credentials)", 0, 0.0

        for attempt in range(max_retries):
            try:
                image_part = Part.from_data(data=meta.image_bytes, mime_type="image/png")
                
                keyword_list = ', '.join(list(self.document_keywords)[:15]) if self.document_keywords else "일반 학습 내용"
                
                prompt = f"""
이 강의의 핵심 주제: {keyword_list}

주변 텍스트: "{meta.adjacent_text}"

이 이미지가 위 주제와 관련있는 **핵심 학습 자료**인지 판단하세요.

✅ KEEP 기준:
- 주제를 구체적으로 설명하는 시각 자료 (차트, 그래프, 다이어그램, 도표, 만화, 사진)
- 주변 텍스트와 긴밀하게 연결된 핵심 콘텐츠

❌ DISCARD 기준:
- 장식용 이미지 (아이콘, 배경, 테두리, 단순 도형)
- 학습 상황 묘사 삽화 (선생님/학생 그림, 공부하는 모습 등) ⚠️ 중요!
- 주제와 무관하거나 일반적인 이미지

⚠️ 주의: "학습 맥락 제공"은 DISCARD입니다. 진짜 교육 콘텐츠만 KEEP하세요.

출력 형식: KEEP 또는 DISCARD로 시작 + 이유 (1-2문장)
"""
                response = self.model.generate_content([image_part, prompt])
                return response.text.strip()
                
            except Exception as e:
                error_msg = str(e)
                
                if "429" in error_msg or "Resource exhausted" in error_msg:
                    if attempt < max_retries - 1:
                        wait_time = (attempt + 1) * 3
                        _log(f"      ⚠️  Rate Limit, {wait_time}초 대기...")
                        time.sleep(wait_time)
                        continue
                    else:
                        return "DISCARD: API rate limit exceeded"
                else:
                    return f"ERROR: {error_msg}"
        
        return "DISCARD: Failed after all retries"

    def run(self, source_path: str):
        """이미지 필터링 실행"""
        from pathlib import Path
        
        file_ext = Path(source_path).suffix.lower()
        _log(f"\n🔍 분석 시작: {os.path.basename(source_path)} ({file_ext})")
        
        if self.auto_extract:
            self.extract_keywords_from_document(source_path)
        
        extractor = UniversalImageExtractor()
        all_meta = extractor.extract(source_path)
        
        _log("\n" + "="*120)
        _log(f"{'Slide':<6} | {'Size':<6} | {'Filter':<12} | {'Result':<12} | {'Reason'}")
        _log("-" * 120)

        final_core = []
        stats = {
            'total': len(all_meta),
            'rule_pass': 0,
            'rule_drop': 0,
            'ai_keep': 0,
            'ai_drop': 0,
        }
        
        for meta in all_meta:
            decision_type, s1_reason = self.step1_rule_check(meta)
            
            final_status = ""
            filter_stage = ""
            detail_reason = ""

            if decision_type == "INCLUDE":
                meta.is_core_content = True
                filter_stage = "1차 (Rule)"
                final_status = "✅ PASS"
                detail_reason = s1_reason
                final_core.append(meta)
                stats['rule_pass'] += 1
                
            elif decision_type == "PENDING":
                filter_stage = "2차 (AI)"
                ai_res = self.step2_gemini_check(meta)
                
                if ai_res.upper().startswith("KEEP"):
                    meta.is_core_content = True
                    final_status = "✅ KEEP"
                    stats['ai_keep'] += 1
                    final_core.append(meta)
                else:
                    final_status = "❌ DROP"
                    stats['ai_drop'] += 1
                    
                detail_reason = ai_res.replace('\n', ' ')
                
            else:
                filter_stage = "1차 (Rule)"
                final_status = "❌ DROP"
                detail_reason = s1_reason
                stats['rule_drop'] += 1

            wrapped_reason = textwrap.wrap(detail_reason, width=70)
            _log(f"{meta.slide_number:<6} | {meta.area_percentage:>5.1f}% | {filter_stage:<12} | {final_status:<12} | {wrapped_reason[0]}")
            for line in wrapped_reason[1:]:
                _log(f"{'':<6} | {'':<6} | {'':<12} | {'':<12} | {line}")
            _log("-" * 120)

        _log("\n" + "="*120)
        _log("📊 최종 결과")
        _log("="*120)
        
        _log(f"\n총 이미지: {stats['total']}개")
        _log(f"\n[1차 필터 - 규칙 기반]")
        _log(f"  ✅ 통과: {stats['rule_pass']}개")
        _log(f"  ❌ 제외: {stats['rule_drop']}개")
        _log(f"  ⚠️  2차 이동: {stats['ai_keep'] + stats['ai_drop']}개")
        
        _log(f"\n[2차 필터 - AI 판단]")
        _log(f"  ✅ 통과: {stats['ai_keep']}개")
        _log(f"  ❌ 제외: {stats['ai_drop']}개")
        
        total_keep = stats['rule_pass'] + stats['ai_keep']
        total_drop = stats['rule_drop'] + stats['ai_drop']
        
        _log(f"\n{'='*120}")
        _log(f"💎 최종 핵심 이미지: {total_keep}개 (1차: {stats['rule_pass']}개 + 2차: {stats['ai_keep']}개)")
        _log(f"🗑️  제외된 이미지: {total_drop}개")
        if stats['total'] > 0:
            _log(f"💰 Vision API 사용: {stats['ai_keep'] + stats['ai_drop']}회 ({(stats['ai_keep'] + stats['ai_drop'])/stats['total']*100:.1f}%)")
        _log(f"{'='*120}\n")
        
        return final_core


if __name__ == "__main__":
    import sys
    
    _log("\n" + "="*120)
    _log("🎯 Improved Hybrid Filter V3 - 이미지 필터링 (pdfplumber)")
    _log("="*120)
    
    if len(sys.argv) > 1:
        source_file = sys.argv[1]
        
        if not os.path.exists(source_file):
            _log(f"\n❌ 파일을 찾을 수 없습니다: {source_file}")
            sys.exit(1)
        
        auto_extract = True
        if len(sys.argv) > 2 and sys.argv[2] in ['--no-auto', '-n']:
            auto_extract = False
            _log("\n⚠️  자동 키워드 추출 비활성화")
        else:
            _log("\n✅ 자동 키워드 추출 활성화")
        
        try:
            pipeline = ImprovedHybridFilterPipeline(auto_extract_keywords=auto_extract)
            core_images = pipeline.run(source_file)
            
            _log(f"\n{'='*120}")
            _log(f"✅ 완료! 핵심 이미지: {len(core_images)}개")
            _log(f"{'='*120}\n")
            
        except Exception as e:
            _log(f"\n❌ 에러 발생: {e}")
            import traceback
            traceback.print_exc()
            sys.exit(1)
    
    else:
        _log("\n사용법:")
        _log("  python improved_hybrid_filter_v3.py <파일경로>")
        _log("\n예시:")
        _log("  python improved_hybrid_filter_v3.py 중등국어1.pdf")
        _log("\n✅ V3 개선사항:")
        _log("  - PyMuPDF (AGPL) → pdfplumber (MIT) 전환")
        _log("  - 라이선스 문제 해결")
        _log("  - OCR 기능 완전 유지 (pypdfium2 + PaddleOCR)")
        _log("  - 텍스트-이미지 중첩 감지 유지")
        _log("  - 색상 복잡도 필터 추가 (텍스트 상자 제거) ⭐")
        _log("  - 기존 v2 기능 모두 유지")
        _log("="*120 + "\n")
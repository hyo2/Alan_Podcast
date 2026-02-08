"""
Document Converter Node
다양한 문서 형식(HWP, DOCX, PPTX, URL, TXT)을 PDF로 변환하는 노드
"""

import os
import tempfile
from pathlib import Path
from typing import Optional, Union
from enum import Enum
import requests
from bs4 import BeautifulSoup
from io import BytesIO

# Document processing
from docx import Document as DocxDocument
from pptx import Presentation

import logging

logger = logging.getLogger(__name__)


class DocumentType(Enum):
    """지원하는 문서 타입"""
    PDF = "pdf"
    HWP = "hwp"
    DOCX = "docx"
    PPTX = "pptx"
    TXT = "txt"      # ✅ 추가
    URL = "url"
    UNKNOWN = "unknown"


class DocumentConverterNode:
    """
    다양한 문서 포맷을 PDF로 변환하는 노드
    
    지원 포맷:
    - PDF: 그대로 사용
    - DOCX: LibreOffice/pypandoc 사용하여 변환
    - PPTX: 각 슬라이드를 이미지로 변환 후 PDF 생성
    - HWP: 외부 변환 서비스 또는 LibreOffice 사용
    - TXT: 텍스트 파일을 PDF로 변환
    - URL: HTML을 추출하여 PDF로 변환
    """
    
    def __init__(self, output_dir: str = "./converted_pdfs"):
        """
        Args:
            output_dir: 변환된 PDF 파일을 저장할 디렉토리
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"DocumentConverterNode initialized. Output dir: {self.output_dir}")
    
    def detect_document_type(self, file_path: str) -> DocumentType:
        """파일 확장자를 기반으로 문서 타입 감지"""
        # URL 체크 (확장자보다 우선)
        # ✅ 백슬래시로 변환된 경우도 감지 (Windows 경로 변환 대응)
        if file_path.startswith(("http://", "https://", "http:\\", "https:\\")):
            # 백슬래시를 슬래시로 복원
            if "\\" in file_path:
                logger.warning(f"⚠️ URL이 백슬래시로 변환됨, 복원: {file_path}")
            return DocumentType.URL
        
        extension = Path(file_path).suffix.lower()
        type_mapping = {
            ".pdf": DocumentType.PDF,
            ".hwp": DocumentType.HWP,
            ".docx": DocumentType.DOCX,
            ".pptx": DocumentType.PPTX,
            ".txt": DocumentType.TXT,  # ✅ 추가
        }
        return type_mapping.get(extension, DocumentType.UNKNOWN)
    
    def convert(self, source: str, output_filename: Optional[str] = None) -> str:
        """
        문서를 PDF로 변환
        
        Args:
            source: 원본 파일 경로 또는 URL
            output_filename: 출력 파일명 (없으면 자동 생성)
        
        Returns:
            변환된 PDF 파일 경로
        """
        doc_type = self.detect_document_type(source)
        logger.info(f"Converting {doc_type.value}: {source}")
        
        if output_filename is None:
            if doc_type == DocumentType.URL:
                output_filename = f"web_content_{hash(source) % 100000}.pdf"
            else:
                output_filename = f"{Path(source).stem}_{doc_type.value}.pdf"
        
        output_path = self.output_dir / output_filename
        
        # 타입별 변환 처리
        conversion_methods = {
            DocumentType.PDF: self._handle_pdf,
            DocumentType.DOCX: self._convert_docx_to_pdf,
            DocumentType.PPTX: self._convert_pptx_to_pdf,
            DocumentType.HWP: self._convert_hwp_to_pdf,
            DocumentType.TXT: self._convert_txt_to_pdf,  # ✅ 추가
            DocumentType.URL: self._convert_url_to_pdf,
        }
        
        if doc_type in conversion_methods:
            return conversion_methods[doc_type](source, str(output_path))
        else:
            raise ValueError(f"Unsupported document type: {doc_type}")
    
    def _wrap_text(self, text: str, max_width: float, canvas_obj) -> list:
        """
        텍스트를 줄바꿈 처리
        
        Args:
            text: 원본 텍스트
            max_width: 최대 너비 (포인트)
            canvas_obj: ReportLab Canvas 객체
        
        Returns:
            줄바꿈된 텍스트 리스트
        """
        words = text.split()
        lines = []
        current_line = ""
        
        for word in words:
            test_line = current_line + " " + word if current_line else word
            
            try:
                # stringWidth로 너비 측정
                if canvas_obj.stringWidth(test_line) < max_width:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
            except:
                # stringWidth 실패 시 글자 수로 대략 판단
                if len(test_line) < 80:
                    current_line = test_line
                else:
                    if current_line:
                        lines.append(current_line)
                    current_line = word
        
        if current_line:
            lines.append(current_line)
        
        # 최소 1줄 반환
        return lines if lines else [text[:100]]
    
    def _handle_pdf(self, source: str, output_path: str) -> str:
        """PDF는 그대로 복사"""
        import shutil
        shutil.copy2(source, output_path)
        logger.info(f"PDF copied to: {output_path}")
        return output_path
    
    def _convert_txt_to_pdf(self, source: str, output_path: str) -> str:
        """
        ✅ TXT 파일을 PDF로 변환
        """
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        
        logger.info(f"Converting TXT to PDF: {source}")
        
        try:
            # 텍스트 파일 읽기
            with open(source, 'r', encoding='utf-8') as f:
                text_content = f.read()
            
            if not text_content.strip():
                raise ValueError("TXT 파일이 비어있습니다")
            
            # PDF 생성
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            # 한글 폰트 등록 시도
            korean_font_registered = False
            try:
                font_paths = [
                    "C:/Windows/Fonts/malgun.ttf",
                    "C:/Windows/Fonts/NanumGothic.ttf",
                    "C:/Windows/Fonts/gulim.ttc",
                    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",  # Linux
                    "/System/Library/Fonts/AppleGothic.ttf",  # macOS
                ]
                for font_path in font_paths:
                    if os.path.exists(font_path):
                        pdfmetrics.registerFont(TTFont('Korean', font_path))
                        korean_font_registered = True
                        logger.info(f"✓ Korean font registered: {font_path}")
                        break
            except Exception as e:
                logger.warning(f"⚠ Font registration failed: {e}")
            
            # 폰트 설정
            if korean_font_registered:
                c.setFont("Korean", 10)
            else:
                c.setFont("Helvetica", 10)
            
            # 제목 (첫 줄 또는 파일명)
            title = text_content.split('\n')[0][:80] if text_content else Path(source).stem
            if korean_font_registered:
                c.setFont("Korean", 14)
            else:
                c.setFont("Helvetica-Bold", 14)
            c.drawString(50, height - 40, title)
            
            # 본문
            if korean_font_registered:
                c.setFont("Korean", 10)
            else:
                c.setFont("Helvetica", 10)
            
            y_position = height - 70
            lines = text_content.split('\n')
            
            for line in lines:
                if not line.strip():
                    y_position -= 14
                    continue
                
                # 긴 줄 자동 줄바꿈
                wrapped_lines = self._wrap_text(line, width - 100, c)
                for wrapped_line in wrapped_lines:
                    if len(wrapped_line) > 120:
                        wrapped_line = wrapped_line[:117] + "..."
                    
                    c.drawString(50, y_position, wrapped_line)
                    y_position -= 14
                    
                    # 페이지 넘김
                    if y_position < 50:
                        c.showPage()
                        if korean_font_registered:
                            c.setFont("Korean", 10)
                        else:
                            c.setFont("Helvetica", 10)
                        y_position = height - 50
            
            c.save()
            logger.info(f"✓ TXT converted to PDF: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"✗ TXT conversion failed: {e}")
            raise
    
    def _convert_docx_to_pdf(self, source: str, output_path: str) -> str:
        """
        DOCX를 PDF로 변환
        LibreOffice 사용 (가장 안정적)
        """
        try:
            # LibreOffice를 사용한 변환
            import subprocess
            
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.output_dir),
                source
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # LibreOffice는 원본 파일명.pdf로 저장
                temp_pdf = self.output_dir / f"{Path(source).stem}.pdf"
                if temp_pdf.exists() and str(temp_pdf) != output_path:
                    temp_pdf.rename(output_path)
                logger.info(f"DOCX converted to: {output_path}")
                return output_path
            else:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")
                
        except Exception as e:
            logger.error(f"DOCX conversion error: {e}")
            # Fallback: 텍스트 추출 후 간단한 PDF 생성
            return self._fallback_docx_conversion(source, output_path)
    
    def _fallback_docx_conversion(self, source: str, output_path: str) -> str:
        """DOCX 변환 실패시 폴백: 텍스트만 추출하여 간단한 PDF 생성"""
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        
        doc = DocxDocument(source)
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        y_position = height - 50
        for para in doc.paragraphs:
            if para.text.strip():
                c.drawString(50, y_position, para.text[:100])  # 간단히 처리
                y_position -= 20
                
                if y_position < 50:
                    c.showPage()
                    y_position = height - 50
        
        c.save()
        logger.info(f"DOCX fallback conversion completed: {output_path}")
        return output_path
    
    def _convert_pptx_to_pdf(self, source: str, output_path: str) -> str:
        """
        PPTX를 PDF로 변환
        각 슬라이드를 이미지로 변환 후 PDF 생성
        """
        try:
            # LibreOffice 사용 (권장)
            import subprocess
            
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.output_dir),
                source
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                temp_pdf = self.output_dir / f"{Path(source).stem}.pdf"
                if temp_pdf.exists() and str(temp_pdf) != output_path:
                    temp_pdf.rename(output_path)
                logger.info(f"PPTX converted to: {output_path}")
                return output_path
            else:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")
                
        except Exception as e:
            logger.error(f"PPTX conversion error: {e}")
            # Fallback: python-pptx로 텍스트만 추출
            return self._fallback_pptx_conversion(source, output_path)
    
    def _fallback_pptx_conversion(self, source: str, output_path: str) -> str:
        """PPTX 폴백: 텍스트 추출 후 PDF 생성"""
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        
        prs = Presentation(source)
        
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4
        
        for slide_num, slide in enumerate(prs.slides):
            c.drawString(50, height - 30, f"Slide {slide_num + 1}")
            y_position = height - 60
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    c.drawString(50, y_position, shape.text[:100])
                    y_position -= 20
                    
                    if y_position < 50:
                        break
            
            c.showPage()
        
        c.save()
        logger.info(f"PPTX fallback conversion completed: {output_path}")
        return output_path
    
    def _convert_hwp_to_pdf(self, source: str, output_path: str) -> str:
        """
        HWP를 PDF로 변환
        
        옵션:
        1. 한글과컴퓨터에서 제공하는 API 사용 (유료)
        2. LibreOffice 사용 (무료, HWP 필터 설치 필요)
        3. 온라인 변환 서비스 사용
        """
        try:
            # LibreOffice로 시도 (HWP 필터가 설치되어 있어야 함)
            import subprocess
            
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.output_dir),
                source
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                temp_pdf = self.output_dir / f"{Path(source).stem}.pdf"
                if temp_pdf.exists() and str(temp_pdf) != output_path:
                    temp_pdf.rename(output_path)
                logger.info(f"HWP converted to: {output_path}")
                return output_path
            else:
                logger.warning(f"LibreOffice HWP conversion failed: {result.stderr}")
                raise Exception("HWP conversion requires LibreOffice with HWP filter")
                
        except Exception as e:
            logger.error(f"HWP conversion error: {e}")
            raise NotImplementedError(
                "HWP conversion requires LibreOffice with HWP filter installed. "
                "Please install: sudo apt-get install libreoffice-writer libreoffice-java-common"
            )
    
    def _convert_url_to_pdf(self, url: str, output_path: str) -> str:
        """
        URL 컨텐츠를 PDF로 변환 (개선 버전)
        - 메인 컨텐츠만 추출
        - 불필요한 요소 제거
        - 한글 폰트 지원
        """
        try:
            # 웹페이지 가져오기
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            response.encoding = 'utf-8'  # UTF-8 인코딩 명시
            
            soup = BeautifulSoup(response.text, 'html.parser')
            
            # 불필요한 요소 제거
            for element in soup(['script', 'style', 'nav', 'header', 'footer', 'aside', 'iframe', 'noscript']):
                element.decompose()
            
            # 메인 컨텐츠 추출 (우선순위 순)
            main_content = None
            
            # 1순위: article, main, role=main
            for selector in ['article', 'main', '[role="main"]', '.content', '#content', '.post-content', '.entry-content', '.document']:
                main_content = soup.select_one(selector)
                if main_content:
                    logger.info(f"✓ Found main content with selector: {selector}")
                    break
            
            # 2순위: body (없으면 전체)
            if not main_content:
                main_content = soup.find('body') or soup
                logger.warning("⚠ Using body as main content")
            
            # 텍스트 추출
            text_content = main_content.get_text(separator='\n', strip=True)
            
            # 빈 줄 정리
            lines = [line.strip() for line in text_content.split('\n') if line.strip()]
            text_content = '\n'.join(lines)
            
            # 최소 기준 완화: 30자 이상이면 시도
            if not text_content or len(text_content) < 30:
                raise ValueError(f"Extracted text too short ({len(text_content)} chars). URL might be inaccessible or have no content.")
            
            # 경고만 표시
            if len(text_content) < 100:
                logger.warning(f"⚠️ URL에서 추출한 텍스트가 짧음 ({len(text_content)}자). 제한적인 정보만 포함될 수 있습니다.")
            
            logger.info(f"✓ Extracted {len(text_content)} characters, {len(lines)} lines")
            
            # PDF 생성
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            
            # 한글 폰트 등록
            korean_font_registered = False
            try:
                font_paths = [
                    "C:/Windows/Fonts/malgun.ttf",
                    "C:/Windows/Fonts/NanumGothic.ttf",
                    "C:/Windows/Fonts/gulim.ttc",
                ]
                for font_path in font_paths:
                    if os.path.exists(font_path):
                        pdfmetrics.registerFont(TTFont('Korean', font_path))
                        korean_font_registered = True
                        logger.info(f"✓ Korean font registered: {font_path}")
                        break
            except Exception as e:
                logger.warning(f"⚠ Font registration failed: {e}")
            
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            # 타이틀
            title = soup.find('title')
            if title:
                title_text = title.get_text().strip()
                if korean_font_registered:
                    c.setFont("Korean", 14)
                else:
                    c.setFont("Helvetica-Bold", 14)
                
                # 타이틀이 너무 길면 자르기
                if len(title_text) > 80:
                    title_text = title_text[:77] + "..."
                c.drawString(50, height - 40, title_text)
            
            # URL 표시
            if korean_font_registered:
                c.setFont("Korean", 9)
            else:
                c.setFont("Helvetica", 9)
            
            url_display = url if len(url) <= 90 else url[:87] + "..."
            c.drawString(50, height - 60, f"Source: {url_display}")
            
            # 구분선
            c.line(50, height - 70, width - 50, height - 70)
            
            # 본문
            if korean_font_registered:
                c.setFont("Korean", 10)
            else:
                c.setFont("Helvetica", 10)
            
            y_position = height - 90
            
            for line in lines[:300]:  # 최대 300줄
                if line.strip():
                    # 긴 줄은 자동으로 나누기
                    wrapped_lines = self._wrap_text(line, width - 100, c)
                    for wrapped_line in wrapped_lines:
                        # 텍스트가 너무 길면 잘라내기
                        if len(wrapped_line) > 120:
                            wrapped_line = wrapped_line[:117] + "..."
                        
                        c.drawString(50, y_position, wrapped_line)
                        y_position -= 14
                        
                        if y_position < 50:
                            c.showPage()
                            if korean_font_registered:
                                c.setFont("Korean", 10)
                            else:
                                c.setFont("Helvetica", 10)
                            y_position = height - 50
            
            c.save()
            logger.info(f"✓ URL converted to PDF: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"✗ URL conversion failed: {e}")
            raise


# 사용 예시
if __name__ == "__main__":
    import sys
    
    logger.warning("이건 보이나?")
    logger.info("\n" + "="*120)
    logger.info("🎯 Document Converter Node")
    logger.info("="*120)
    
    # CLI 인자가 있으면 그것 사용
    if len(sys.argv) > 1:
        source_file = sys.argv[1]
        
        if not source_file.startswith('http'):
            if not os.path.exists(source_file):
                logger.error(f"\n❌ 파일을 찾을 수 없습니다: {source_file}")
                sys.exit(1)
        
        converter = DocumentConverterNode(output_dir="./test_output")
        
        try:
            output = converter.convert(source_file)
            logger.info(f"\n✓ 변환 완료: {output}")
        except Exception as e:
            logger.error(f"\n✗ 변환 실패: {e}")
            import traceback
            traceback.print_exc()
            sys.exit(1)
    
    else:
        # CLI 인자 없으면 사용법 표시
        logger.info("\n사용법:")
        logger.info("  python document_converter_node.py <파일경로 또는 URL>")
        logger.info("\n✅ 지원 형식:")
        logger.info("  - PDF (원본 복사)")
        logger.info("  - DOCX (LibreOffice 변환)")
        logger.info("  - PPTX (LibreOffice 변환)")
        logger.info("  - TXT (텍스트 → PDF)")
        logger.info("  - URL (웹페이지 크롤링)")
        logger.info("\n예제:")
        logger.info("  python document_converter_node.py sample.docx")
        logger.info("  python document_converter_node.py notes.txt")
        logger.info("  python document_converter_node.py https://example.com")
        logger.info("="*120 + "\n")